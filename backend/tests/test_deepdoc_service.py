"""DeepDoc routing regressions."""

from io import BytesIO
from zipfile import ZIP_DEFLATED, ZIP_STORED, ZipFile

from app.services.deepdoc_service import parse_document


def test_text_parser_remains_available_with_lazy_parser_exports():
    blocks = parse_document("notes.txt", "第一段\nsecond section".encode())

    assert blocks
    assert all(block["content_type"] == "text" for block in blocks)
    assert "第一段" in "\n".join(block["content"] for block in blocks)


def test_parser_package_public_contract():
    from deepdoc.parser import JsonParser, MarkdownParser, TxtParser

    assert TxtParser.__name__ == "RAGFlowTxtParser"
    assert JsonParser.__name__ == "RAGFlowJsonParser"
    assert MarkdownParser.__name__ == "RAGFlowMarkdownParser"


def test_docx_parser_extracts_generated_document():
    from docx import Document

    output = BytesIO()
    document = Document()
    document.add_heading("ORBIT-DOCX", level=1)
    document.add_paragraph("Generated Word ingestion evidence")
    document.save(output)

    blocks = parse_document("evidence.docx", output.getvalue())

    assert "ORBIT-DOCX" in "\n".join(block["content"] for block in blocks)


def test_xlsx_parser_extracts_generated_workbook():
    from openpyxl import Workbook

    output = BytesIO()
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Evidence"
    worksheet.append(["marker", "status"])
    worksheet.append(["ORBIT-XLSX", "ready"])
    workbook.save(output)

    blocks = parse_document("evidence.xlsx", output.getvalue())

    assert "ORBIT-XLSX" in "\n".join(block["content"] for block in blocks)


def test_pptx_parser_extracts_generated_single_slide_with_page_number():
    from pptx import Presentation

    output = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "ORBIT-PPTX"
    slide.placeholders[1].text = "Generated presentation ingestion evidence"
    presentation.save(output)

    blocks = parse_document("evidence.pptx", output.getvalue())

    assert len(blocks) == 1
    assert blocks[0]["page"] == 0
    assert "ORBIT-PPTX" in blocks[0]["content"]


def test_markdown_parser_extracts_text_and_table():
    source = b"# ORBIT-MD\n\n| marker | status |\n| --- | --- |\n| table | ready |\n"

    blocks = parse_document("evidence.md", source)

    assert "ORBIT-MD" in "\n".join(block["content"] for block in blocks)
    assert {block["content_type"] for block in blocks} == {"text", "table"}


def test_html_parser_falls_back_when_nltk_data_is_missing(monkeypatch):
    import infinity.rag_tokenizer

    def missing_nltk_data(*args, **kwargs):
        raise LookupError("punkt_tab missing")

    monkeypatch.setattr(infinity.rag_tokenizer.RagTokenizer, "tokenize", missing_nltk_data)
    blocks = parse_document(
        "evidence.html",
        b"<html><body><h1>ORBIT-HTML</h1><p>Offline evidence</p></body></html>",
    )

    assert "ORBIT-HTML" in "\n".join(block["content"] for block in blocks)


def test_epub_parser_falls_back_when_nltk_data_is_missing(monkeypatch):
    import infinity.rag_tokenizer

    def missing_nltk_data(*args, **kwargs):
        raise LookupError("punkt_tab missing")

    monkeypatch.setattr(infinity.rag_tokenizer.RagTokenizer, "tokenize", missing_nltk_data)
    output = BytesIO()
    with ZipFile(output, "w") as epub:
        epub.writestr("mimetype", "application/epub+zip", compress_type=ZIP_STORED)
        epub.writestr(
            "chapter.xhtml",
            '<html xmlns="http://www.w3.org/1999/xhtml"><body><h1>ORBIT-EPUB</h1></body></html>',
            compress_type=ZIP_DEFLATED,
        )

    blocks = parse_document("evidence.epub", output.getvalue())

    assert "ORBIT-EPUB" in "\n".join(block["content"] for block in blocks)
